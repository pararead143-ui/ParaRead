from rest_framework.decorators import api_view, permission_classes
from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response
from rest_framework import status
from .models import Material, Vocabulary
from .serializers import MaterialSerializer, VocabularySerializer
from .segmentation import segment_text
from .summarization import summarize_text
from .quiz import generate_quiz_from_summary
from .education import generate_educational_insights

import re
import pdfplumber
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
from .models import Vocabulary
import requests


# Make sure Tesseract is correctly pointed
pytesseract.pytesseract.tesseract_cmd = r"C:\Users\Jamie\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"


def clean_text(text: str) -> str:
    # Remove emojis
    text = re.sub(
        r'[\U0001F600-\U0001F64F\U0001F300-\U0001F5FF\U0001F680-\U0001F6FF\U0001F1E0-\U0001F1FF]',
        '',
        text
    )
    # Remove short lines (less than 5 words)
    lines = [line for line in text.split('\n') if len(line.split()) > 4]
    return '\n'.join(lines)


# --- SEGMENTATION ---
@api_view(["POST"])
@permission_classes([IsAuthenticated])
def segment_view(request):
    text = request.data.get("text", "")
    title = request.data.get("title", "")
    material_id = request.data.get("material_id")  # Optional: existing material

    if not text.strip():
        return Response({"error": "No text provided"}, status=status.HTTP_400_BAD_REQUEST)

    segmented = segment_text(text)

    if material_id:
        # ✅ Update existing material WITHOUT wiping summary_data
        try:
            material = Material.objects.get(id=material_id, user=request.user)

            material.segmented_data = segmented
            # Optional: keep raw_text/title in sync with latest text
            material.raw_text = text
            if title:
                material.title = title

            fields_to_update = ["segmented_data", "raw_text"]
            if title:
                fields_to_update.append("title")

            # ✅ Only update these fields; keeps summary_data intact
            material.save(update_fields=fields_to_update)

        except Material.DoesNotExist:
            # If material_id is invalid, create a new one
            material = Material.objects.create(
                user=request.user,
                title=title or "Untitled",
                raw_text=text,
                segmented_data=segmented
            )
    else:
        # No material_id, create new
        material = Material.objects.create(
            user=request.user,
            title=title or "Untitled",
            raw_text=text,
            segmented_data=segmented
        )

    return Response(MaterialSerializer(material).data, status=status.HTTP_200_OK)


# --- SUMMARIZATION ---
@api_view(["POST"])
@permission_classes([IsAuthenticated])
def summarize_view(request):
    text = request.data.get("text", "")
    material_id = request.data.get("material_id")
    title = request.data.get("title", "")

    if not text.strip():
        return Response({"error": "No text provided"}, status=status.HTTP_400_BAD_REQUEST)

    # Generate the summary
    summary_obj = summarize_text(text)
    summary_text = summary_obj.get("summary") if isinstance(summary_obj, dict) else str(summary_obj)

    if material_id:
        try:
            # ✅ Update existing material WITHOUT wiping segmented_data
            material = Material.objects.get(id=material_id, user=request.user)

            material.summary_data = {"summary": summary_text}
            # Optional: keep raw_text/title in sync with latest text
            material.raw_text = text
            if title:
                material.title = title

            fields_to_update = ["summary_data", "raw_text"]
            if title:
                fields_to_update.append("title")

            # ✅ Only update these fields; keeps segmented_data intact
            material.save(update_fields=fields_to_update)

        except Material.DoesNotExist:
            # Invalid ID: create new material
            material = Material.objects.create(
                user=request.user,
                raw_text=text,
                summary_data={"summary": summary_text},
                title=title or "Untitled"
            )
    else:
        # No material_id: create new material
        material = Material.objects.create(
            user=request.user,
            raw_text=text,
            summary_data={"summary": summary_text},
            title=title or "Untitled"
        )

    # Return the summary and material ID to frontend
    return Response(
        {"summary": summary_text, "id": material.id},
        status=status.HTTP_200_OK
    )


# --- GET ALL MATERIALS ---
@api_view(["GET"])
@permission_classes([IsAuthenticated])
def materials_list(request):
    materials = Material.objects.filter(user=request.user).order_by("-created_at")
    serializer = MaterialSerializer(materials, many=True)
    return Response(serializer.data, status=status.HTTP_200_OK)


@api_view(["GET", "DELETE", "PATCH"])
@permission_classes([IsAuthenticated])
def material_detail(request, material_id):
    try:
        material = Material.objects.get(id=material_id, user=request.user)
    except Material.DoesNotExist:
        return Response({"error": "Material not found"}, status=status.HTTP_404_NOT_FOUND)

    if request.method == "GET":
        serializer = MaterialSerializer(material)
        return Response(serializer.data, status=status.HTTP_200_OK)

    if request.method == "PATCH":
        serializer = MaterialSerializer(material, data=request.data, partial=True)
        if serializer.is_valid():
            serializer.save()
            return Response(serializer.data, status=status.HTTP_200_OK)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

    # DELETE
    material.delete()
    return Response({"message": "Material deleted successfully."}, status=status.HTTP_204_NO_CONTENT)

# --- EDUCATIONAL INSIGHTS ---
@api_view(["GET"])
@permission_classes([IsAuthenticated])
def educational_insights(request, material_id):
    try:
        material = Material.objects.get(id=material_id, user=request.user)
    except Material.DoesNotExist:
        return Response({"error": "Material not found"}, status=status.HTTP_404_NOT_FOUND)

    if not material.segmented_data:
        return Response({"error": "Material has no segmentation yet"}, status=status.HTTP_400_BAD_REQUEST)

    insights = generate_educational_insights(material.segmented_data)
    return Response(insights, status=status.HTTP_200_OK)


# --- FILE UPLOAD (Does NOT create Material) ---
@api_view(["POST"])
@permission_classes([IsAuthenticated])
def upload_file(request):
    if request.FILES.get('file'):
        file = request.FILES['file']
        ext = file.name.split('.')[-1].lower()
        text = ""

        if ext == "pdf":
            with pdfplumber.open(file) as pdf:
                text = "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())
        elif ext == "docx":
            doc = docx.Document(file)
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext == "txt":
            text = file.read().decode('utf-8')
        elif ext == "pptx":
            prs = Presentation(file)
            text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
        elif ext in ["jpg", "png"]:
            img = Image.open(file).convert("RGB")
            text = pytesseract.image_to_string(img)
        else:
            return Response({"error": "Unsupported file type"}, status=status.HTTP_400_BAD_REQUEST)

        cleaned = clean_text(text)

        # DO NOT create Material here, just return cleaned text
        return Response({"cleaned_text": cleaned}, status=status.HTTP_200_OK)

    return Response({"error": "No file uploaded"}, status=status.HTTP_400_BAD_REQUEST)


# --- QUIZ GENERATION ---
@api_view(["POST"])
@permission_classes([IsAuthenticated])
def generate_quiz(request, material_id):
    try:
        material = Material.objects.get(id=material_id, user=request.user)
    except Material.DoesNotExist:
        return Response({"error": "Material not found"}, status=status.HTTP_404_NOT_FOUND)

    summary_text = material.summary_data.get("summary", "") if isinstance(material.summary_data, dict) else ""

    if not summary_text.strip():
        return Response({"error": "No summary available to generate a quiz."}, status=status.HTTP_400_BAD_REQUEST)

    quiz_result = generate_quiz_from_summary(summary_text, num_questions=5)

    if not quiz_result:
        return Response({"error": "Quiz could not be generated from the summary."}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    material.quiz_data = quiz_result
    # ✅ Optional: also protect other fields by updating only quiz_data
    material.save(update_fields=["quiz_data"])

    return Response({"quiz": quiz_result}, status=status.HTTP_200_OK)


@api_view(["GET", "POST"])
@permission_classes([IsAuthenticated])
def vocabulary_list_create(request):
    if request.method == "GET":
        vocab = Vocabulary.objects.filter(user=request.user)
        serializer = VocabularySerializer(vocab, many=True)
        return Response(serializer.data, status=status.HTTP_200_OK)

    # POST
    word = request.data.get("word", "").strip()
    meaning = request.data.get("meaning", "").strip()  # frontend sends 'meaning'

    if not word:
        return Response({"error": "Word is required."}, status=status.HTTP_400_BAD_REQUEST)

    vocab = Vocabulary.objects.create(
        user=request.user,
        word=word,
        meaning=meaning
    )
    serializer = VocabularySerializer(vocab)
    return Response(serializer.data, status=status.HTTP_201_CREATED)


@api_view(["DELETE"])
@permission_classes([IsAuthenticated])
def vocabulary_delete(request, pk):
    """Delete a saved vocabulary word"""
    try:
        vocab = Vocabulary.objects.get(id=pk, user=request.user)
    except Vocabulary.DoesNotExist:
        return Response({"error": "Word not found."}, status=status.HTTP_404_NOT_FOUND)

    vocab.delete()
    return Response({"message": "Word deleted successfully."}, status=status.HTTP_204_NO_CONTENT)

@api_view(["GET"])
@permission_classes([IsAuthenticated])
def lookup_definition(request):
    word = (request.query_params.get("word") or "").strip()
    if not word:
        return Response({"error": "Missing word"}, status=status.HTTP_400_BAD_REQUEST)

    # 1) Check if already saved in DB
    existing = Vocabulary.objects.filter(user=request.user, word__iexact=word).first()
    if existing:
        return Response({
            "word": existing.word,
            "meaning": existing.meaning,
            "example": existing.example
        })

    # 2) Fetch from Free Dictionary API
    try:
        url = f"https://api.dictionaryapi.dev/api/v2/entries/en/{word}"
        r = requests.get(url, timeout=10)
        if r.status_code != 200:
            return Response({"error": "No definition found"}, status=404)

        data = r.json()

        # Pick the first definition
        meaning = ""
        example = ""
        try:
            meaning = data[0]["meanings"][0]["definitions"][0].get("definition", "")
            example = data[0]["meanings"][0]["definitions"][0].get("example", "")
        except Exception:
            pass

        if not meaning:
            return Response({"error": "No definition found"}, status=404)

        # 3) Save into your Vocabulary
        vocab = Vocabulary.objects.create(
            user=request.user,
            word=word,
            meaning=meaning,
            example=example or ""
        )

        return Response({
            "word": vocab.word,
            "meaning": vocab.meaning,
            "example": vocab.example
        })

    except requests.RequestException:
        return Response({"error": "Dictionary service unavailable"}, status=503)
